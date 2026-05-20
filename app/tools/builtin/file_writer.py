"""
【模块说明】文件写入工具（FileWriterTool）— 让 AI 能把生成的内容保存为文件

当 AI 生成了报告、代码、表格等内容，用户想要保存到文件时，
这个工具负责把内容写入服务器上用户专属的 generated/ 目录。

【文件保存路径】
  {上传根目录}/{用户ID}/generated/{可选子目录}/{文件名}

【支持的格式】
  文本和代码   — .txt .md .py .js .ts .html .xml .log .svg 等
  数据文件     — .json .yaml .csv .tsv 等
  Office 文档  — .docx .xlsx .pptx（使用 python-docx/openpyxl 库生成）
  二进制文件   — .pdf .png .jpg 等（需传入 base64 编码内容）

【注意】这是需要用户授权的操作，调用前会触发授权检查。

文件保存工具 — 将各种格式内容写入服务端用户生成文件目录。

保存路径规则：{UPLOAD_ROOT}/{user_id}/generated/{subdir?}/{filename}
  UPLOAD_ROOT 由 config/system_config.yaml file_storage.upload_dir 指定（默认 data/uploads）
  user_id 取自 context["user_id"]

支持格式：
  文本类   : .txt .md .rst .html .xml .log .svg 以及所有代码扩展名
  结构化数据: .json .jsonl .yaml .yml .toml .ini .env .cfg .conf
  表格     : .csv .tsv
  Office   : .docx .xlsx .pptx
  二进制   : .pdf .png .jpg .jpeg .gif .webp .bmp（base64 输入）

内容传入约定：
  - 文本/代码/INI/TOML    → 字符串
  - JSON / YAML          → 字符串 或 dict/list（自动序列化）
  - JSONL                → 字符串 或 list（每项序列化为一行）
  - CSV / TSV            → 字符串 或 list[list] / list[dict]（自动转行）
  - Excel (.xlsx)        → list[list]（第一行为表头）或 list[dict]
  - Word (.docx)         → 字符串（\\n\\n 分段）或 list[str]（每项为一段落）
  - PowerPoint (.pptx)   → list[dict]，每项含 "title" 和 "content"（str 或 list[str]）
  - PDF / 图片           → base64 编码字符串
"""

from __future__ import annotations

import base64
import csv
import json
import logging
import os
import sys
from datetime import datetime
from io import StringIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from app.utils.paths import PROJECT_ROOT
from app.core.file_storage import UPLOAD_ROOT, MAX_FILE_SIZE, SUBDIR_GENERATED
from app.tools.base import BaseTool, EXEC_SERVER, VIS_PUBLIC

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════════
# 软导入辅助
# ══════════════════════════════════════════════════════════════════

def _try_import(module: str, package: Optional[str] = None):
    try:
        import importlib
        return importlib.import_module(module)
    except ImportError:
        pkg = package or module
        logger.debug("可选依赖 %s 未安装（pip install %s）", module, pkg)
        return None


# ══════════════════════════════════════════════════════════════════
# 格式分类
# ══════════════════════════════════════════════════════════════════

_EXT_CATEGORY = {
    # 纯文本
    ".txt": "text", ".md": "text", ".rst": "text", ".log": "text",
    ".html": "text", ".htm": "text", ".xml": "text", ".svg": "text",
    # 结构化文本
    ".json": "json", ".jsonl": "jsonl",
    ".yaml": "yaml", ".yml": "yaml",
    ".toml": "toml",
    ".ini": "text", ".env": "text", ".cfg": "text", ".conf": "text",
    # 表格
    ".csv": "csv", ".tsv": "tsv",
    # Office
    ".docx": "word",
    ".xlsx": "excel",
    ".pptx": "ppt",
    # 二进制
    ".pdf": "binary",
    ".png": "binary", ".jpg": "binary", ".jpeg": "binary",
    ".gif": "binary", ".webp": "binary", ".bmp": "binary",
    ".tiff": "binary", ".tif": "binary",
    # 代码（文本处理）
    ".py": "text", ".js": "text", ".ts": "text", ".go": "text",
    ".rs": "text", ".java": "text", ".cpp": "text", ".c": "text",
    ".h": "text", ".sh": "text", ".bash": "text", ".sql": "text",
    ".r": "text", ".rb": "text", ".php": "text", ".kt": "text",
    ".swift": "text", ".css": "text", ".scss": "text", ".less": "text",
}


# ══════════════════════════════════════════════════════════════════
# 路径解析与安全校验
# ══════════════════════════════════════════════════════════════════

def _resolve_save_path(filename: str, user_id: str, subdir: Optional[str]) -> Path:
    """
    解析保存路径，确保结果在用户 generated/ 沙箱内。

    最终路径：UPLOAD_ROOT / user_id / generated / [subdir] / filename
    """
    # generated/ 沙箱边界
    gen_root = UPLOAD_ROOT / user_id / SUBDIR_GENERATED

    if subdir:
        clean_subdir = Path(os.path.normpath(subdir))
        if clean_subdir.is_absolute():
            raise PermissionError(f"subdir 不能为绝对路径：{subdir}")
        target_dir = gen_root / clean_subdir
    else:
        target_dir = gen_root

    target = Path(os.path.normpath(target_dir / filename))

    try:
        target.relative_to(Path(os.path.normpath(gen_root)))
    except ValueError:
        raise PermissionError(
            f"文件名 '{filename}' 超出 generated 目录范围（检测到路径穿越攻击）"
        )
    return target


# ══════════════════════════════════════════════════════════════════
# 各格式写入器
# ══════════════════════════════════════════════════════════════════

def _write_text(
    path: Path,
    content: Any,
    encoding: str,
) -> Dict[str, Any]:
    """写入纯文本/代码文件。"""
    if not isinstance(content, str):
        content = str(content)
    data = content.encode(encoding, errors="replace")
    path.write_bytes(data)
    return {"chars": len(content), "encoding": encoding}


def _write_json(path: Path, content: Any, encoding: str) -> Dict[str, Any]:
    """写入 JSON 文件。dict/list 自动序列化。"""
    if isinstance(content, (dict, list)):
        text = json.dumps(content, ensure_ascii=False, indent=2)
    elif isinstance(content, str):
        # 验证 JSON 合法性
        try:
            json.loads(content)
        except json.JSONDecodeError as e:
            raise ValueError(f"content 不是合法的 JSON 字符串: {e}")
        text = content
    else:
        text = json.dumps(content, ensure_ascii=False, indent=2)
    path.write_bytes(text.encode(encoding, errors="replace"))
    return {"chars": len(text)}


def _write_jsonl(path: Path, content: Any, encoding: str) -> Dict[str, Any]:
    """写入 JSONL 文件。list 中每项序列化为一行。"""
    if isinstance(content, list):
        lines = [json.dumps(item, ensure_ascii=False) for item in content]
        text = "\n".join(lines) + "\n"
    elif isinstance(content, str):
        text = content
    else:
        text = json.dumps(content, ensure_ascii=False) + "\n"
    path.write_bytes(text.encode(encoding, errors="replace"))
    return {"lines": text.count("\n")}


def _write_yaml(path: Path, content: Any, encoding: str) -> Dict[str, Any]:
    """写入 YAML 文件。dict/list 自动序列化。"""
    yaml = _try_import("yaml")
    if isinstance(content, (dict, list)):
        if not yaml:
            raise ImportError("写入 YAML 需要安装 PyYAML（pip install pyyaml）")
        text = yaml.dump(content, allow_unicode=True, default_flow_style=False)
    elif isinstance(content, str):
        text = content
    else:
        text = str(content)
    path.write_bytes(text.encode(encoding, errors="replace"))
    return {"chars": len(text)}


def _write_toml(path: Path, content: Any, encoding: str) -> Dict[str, Any]:
    """写入 TOML 文件。dict 时优先用 tomlkit，否则降级为字符串写入。"""
    if isinstance(content, dict):
        tomlkit = _try_import("tomlkit")
        if tomlkit:
            text = tomlkit.dumps(content)
        else:
            # 无 tomlkit 时退化为 JSON 兜底写入（内容合法但非 TOML 原生格式）
            logger.warning("tomlkit 未安装，TOML dict 内容以 JSON 形式写入（pip install tomlkit）")
            text = json.dumps(content, ensure_ascii=False, indent=2)
    elif isinstance(content, str):
        text = content
    else:
        text = str(content)
    path.write_bytes(text.encode(encoding, errors="replace"))
    return {"chars": len(text)}


def _write_csv(
    path: Path,
    content: Any,
    delimiter: str,
    encoding: str,
) -> Dict[str, Any]:
    """写入 CSV/TSV 文件。支持字符串、list[list]、list[dict] 三种输入。"""
    if isinstance(content, str):
        path.write_bytes(content.encode(encoding, errors="replace"))
        return {"rows": content.count("\n")}

    buf = StringIO()
    if isinstance(content, list) and content:
        if isinstance(content[0], dict):
            headers = list(content[0].keys())
            writer = csv.DictWriter(buf, fieldnames=headers, delimiter=delimiter,
                                    lineterminator="\n")
            writer.writeheader()
            writer.writerows(content)
            rows = len(content)
        else:
            writer = csv.writer(buf, delimiter=delimiter, lineterminator="\n")
            writer.writerows(content)
            rows = len(content)
    else:
        raise ValueError("CSV/TSV content 必须为字符串、list[list] 或 list[dict]")

    path.write_bytes(buf.getvalue().encode(encoding, errors="replace"))
    return {"rows": rows}


def _write_excel(
    path: Path,
    content: Any,
    sheet_name: str,
) -> Dict[str, Any]:
    """写入 Excel (.xlsx) 文件。content 为 list[list] 或 list[dict]。"""
    openpyxl = _try_import("openpyxl")
    if not openpyxl:
        raise ImportError("写入 Excel 需要安装 openpyxl（pip install openpyxl）")

    if not isinstance(content, list):
        raise ValueError("Excel content 必须为 list[list] 或 list[dict]")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name

    if content and isinstance(content[0], dict):
        headers = list(content[0].keys())
        ws.append(headers)
        for row in content:
            ws.append([row.get(h, "") for h in headers])
        rows = len(content)
    elif content and isinstance(content[0], (list, tuple)):
        for row in content:
            ws.append(list(row))
        rows = len(content)
    else:
        raise ValueError("Excel content 的每行必须为 list 或 dict")

    wb.save(path)
    return {"rows": rows, "sheet": sheet_name}


def _write_word(path: Path, content: Any) -> Dict[str, Any]:
    """写入 Word (.docx) 文件。content 为字符串或 list[str]（段落）。"""
    docx_mod = _try_import("docx", package="python-docx")
    if not docx_mod:
        raise ImportError("写入 Word 需要安装 python-docx（pip install python-docx）")

    doc = docx_mod.Document()
    paragraphs = 0

    if isinstance(content, list):
        items = content
    elif isinstance(content, str):
        items = content.split("\n\n")
    else:
        items = [str(content)]

    for item in items:
        text = item.strip() if isinstance(item, str) else str(item).strip()
        if not text:
            continue
        doc.add_paragraph(text)
        paragraphs += 1

    doc.save(path)
    return {"paragraphs": paragraphs}


def _write_ppt(path: Path, content: Any) -> Dict[str, Any]:
    """
    写入 PowerPoint (.pptx) 文件。

    content 格式（list[dict]）：
      [{"title": "标题", "content": "正文" 或 ["要点1", "要点2"]}, ...]

    也可传字符串（每行作为一张幻灯片的正文，以 --- 分隔幻灯片）。
    """
    pptx_mod = _try_import("pptx", package="python-pptx")
    if not pptx_mod:
        raise ImportError("写入 PowerPoint 需要安装 python-pptx（pip install python-pptx）")

    from pptx.util import Inches, Pt

    prs = pptx_mod.Presentation()
    slide_layout = prs.slide_layouts[1]  # 标题 + 内容布局
    slides_written = 0

    if isinstance(content, str):
        # 按 --- 分割幻灯片
        slide_blocks = [b.strip() for b in content.split("---") if b.strip()]
        content = [{"title": "", "content": b} for b in slide_blocks]

    if not isinstance(content, list):
        content = [{"title": "", "content": str(content)}]

    for slide_data in content:
        if not isinstance(slide_data, dict):
            slide_data = {"title": "", "content": str(slide_data)}

        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_text = slide_data.get("title", "")
        if title_shape and title_text:
            title_shape.text = str(title_text)

        body_content = slide_data.get("content", "")
        tf = body_shape.text_frame
        tf.clear()
        if isinstance(body_content, list):
            for i, item in enumerate(body_content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
                p.level = 0
        else:
            tf.paragraphs[0].text = str(body_content)

        slides_written += 1

    prs.save(path)
    return {"slides": slides_written}


def _write_binary(path: Path, content: Any) -> Dict[str, Any]:
    """写入二进制文件（PDF、图片等）。content 必须为 base64 编码字符串或 bytes。"""
    if isinstance(content, bytes):
        raw = content
    elif isinstance(content, str):
        # 去掉可能存在的 Data URL 前缀（如 data:image/png;base64,xxx）
        if "base64," in content:
            content = content.split("base64,", 1)[1]
        try:
            raw = base64.b64decode(content)
        except Exception as e:
            raise ValueError(f"base64 解码失败: {e}")
    else:
        raise ValueError("二进制文件 content 必须为 base64 字符串或 bytes")
    path.write_bytes(raw)
    return {"size_bytes": len(raw)}


# ══════════════════════════════════════════════════════════════════
# 主工具类
# ══════════════════════════════════════════════════════════════════

class FileWriterTool(BaseTool):
    """
    文件保存工具，将各种格式内容写入服务端用户隔离目录。

    保存路径：{upload_dir}/{username}/{subdir?}/{filename}
    upload_dir 由 config/system_config.yaml file_storage.upload_dir 配置。
    """

    name          = "file_writer"
    description   = (
        "将内容保存为指定格式的文件，存储到服务端用户目录。"
        "支持文本、JSON/YAML、CSV/TSV、Excel、Word、PowerPoint 及二进制（PDF/图片）格式。"
        "返回保存后的相对路径及文件元信息。"
    )
    exec_location = EXEC_SERVER
    visibility    = VIS_PUBLIC
    dangerous_ops = ["write"]

    parameters_schema = {
        "filename": {
            "type":        "string",
            "description": (
                "保存的文件名（含扩展名）。支持单级子路径（如 reports/summary.xlsx），"
                "但不允许 ../ 路径穿越。"
            ),
            "required": True,
        },
        "content": {
            "type":        ["string", "object", "array"],
            "description": (
                "文件内容。"
                "文本/代码/INI/TOML：字符串；"
                "JSON/YAML：字符串或 dict/list；"
                "JSONL：字符串或 list；"
                "CSV/TSV：字符串、list[list] 或 list[dict]；"
                "Excel：list[list]（第一行为表头）或 list[dict]；"
                "Word：字符串（\\n\\n 分段）或 list[str]；"
                "PowerPoint：list[dict]，每项含 title 和 content 字段；"
                "PDF/图片：base64 编码字符串。"
            ),
            "required": True,
        },
        "overwrite": {
            "type":        "boolean",
            "description": "文件已存在时是否覆盖（默认 false）。",
            "default":     False,
        },
        "subdir": {
            "type":        "string",
            "description": "在用户目录下创建的子目录（如 'documents/2026'），可选。",
            "required":    False,
        },
        "encoding": {
            "type":        "string",
            "description": "文本文件编码（默认 utf-8）。",
            "default":     "utf-8",
            "required":    False,
        },
        "sheet_name": {
            "type":        "string",
            "description": "Excel 工作表名称（默认 Sheet1）。",
            "default":     "Sheet1",
            "required":    False,
        },
    }

    async def execute(self, params: Dict[str, Any], context: Dict[str, Any]) -> Dict[str, Any]:
        filename   = params["filename"]
        content    = params["content"]
        overwrite  = bool(params.get("overwrite", False))
        subdir     = params.get("subdir") or None
        encoding   = params.get("encoding") or "utf-8"
        sheet_name = params.get("sheet_name") or "Sheet1"

        user_id = context.get("user_id") or "anonymous"

        # ── 路径解析与安全校验 ──────────────────────────────────────
        try:
            save_path = _resolve_save_path(filename, user_id, subdir)
        except PermissionError as e:
            return {"result": str(e), "success": False, "metadata": {}}

        # ── 已存在文件检查 ─────────────────────────────────────────
        if save_path.exists() and not overwrite:
            return {
                "result": f"文件已存在：{save_path.name}（设置 overwrite=true 可覆盖）",
                "success": False,
                "metadata": {"path": str(save_path.relative_to(PROJECT_ROOT))},
            }

        # ── 文件大小预检（仅对字符串/bytes 类型内容有效）──────────────
        if isinstance(content, (str, bytes)):
            estimated = len(content.encode() if isinstance(content, str) else content)
            if estimated > MAX_FILE_SIZE:
                mb = MAX_FILE_SIZE // (1024 * 1024)
                return {
                    "result": f"内容超出最大文件限制 {mb} MB",
                    "success": False,
                    "metadata": {},
                }

        # ── 创建目录 ───────────────────────────────────────────────
        try:
            save_path.parent.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            return {
                "result": f"创建目录失败：{e}",
                "success": False,
                "metadata": {},
            }

        # ── 按格式分派写入 ──────────────────────────────────────────
        ext      = save_path.suffix.lower()
        category = _EXT_CATEGORY.get(ext, "text")
        fmt_meta: Dict[str, Any] = {}

        try:
            if category == "json":
                fmt_meta = _write_json(save_path, content, encoding)

            elif category == "jsonl":
                fmt_meta = _write_jsonl(save_path, content, encoding)

            elif category == "yaml":
                fmt_meta = _write_yaml(save_path, content, encoding)

            elif category == "toml":
                fmt_meta = _write_toml(save_path, content, encoding)

            elif category == "csv":
                fmt_meta = _write_csv(save_path, content, ",", encoding)

            elif category == "tsv":
                fmt_meta = _write_csv(save_path, content, "\t", encoding)

            elif category == "excel":
                fmt_meta = _write_excel(save_path, content, sheet_name)

            elif category == "word":
                fmt_meta = _write_word(save_path, content)

            elif category == "ppt":
                fmt_meta = _write_ppt(save_path, content)

            elif category == "binary":
                fmt_meta = _write_binary(save_path, content)

            else:
                # 文本/代码及一切未知扩展名
                fmt_meta = _write_text(save_path, content, encoding)

        except (ImportError, ValueError) as e:
            return {"result": str(e), "success": False, "metadata": {}}
        except Exception as e:
            logger.exception("文件写入工具异常 path=%s", save_path)
            return {
                "result":  f"文件写入时发生异常：{e}",
                "success": False,
                "metadata": {},
            }

        # ── 构建返回元信息 ─────────────────────────────────────────
        stat = save_path.stat()
        rel_path = str(save_path.relative_to(PROJECT_ROOT))
        metadata = {
            "filename":     save_path.name,
            "saved_path":   rel_path,
            "size_bytes":   stat.st_size,
            "extension":    ext,
            "file_type":    category,
            "saved_at":     datetime.now().isoformat(),
            **fmt_meta,
        }

        logger.info(
            "文件写入成功: user=%s path=%s size=%d bytes",
            user_id, rel_path, stat.st_size,
        )
        return {
            "result": {
                "saved_path": rel_path,
                "metadata":   metadata,
            },
            "success":  True,
            "metadata": {"tool": self.name, "saved_path": rel_path},
        }


# ── 自动注册 ──────────────────────────────────────────────────────

def _register() -> None:
    try:
        from app.tools.registry import registry
        tool_instance = FileWriterTool()
        if not registry.get(tool_instance.name):
            registry.register(tool_instance)
            logger.debug("已注册内置工具: %s", tool_instance.name)
    except Exception as e:
        logger.warning("注册 FileWriterTool 失败: %s", e)


_register()
